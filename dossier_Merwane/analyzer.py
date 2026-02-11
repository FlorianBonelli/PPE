# analyzer.py
import os
import pandas as pd
import json
import csv
import docx
import mimetypes
import zipfile
import tempfile
import openpyxl 
import extract_msg  
from pptx import Presentation
from utils_pdf import lire_pdf_robuste, lire_image_directe

MAX_SIZE_BYTES = 50 * 1024 * 1024

def detect_mime(file_path):
    mime, _ = mimetypes.guess_type(file_path)
    return mime or "application/octet-stream"

def parse_pdf_robuste_wrapper(path):
    try:
        with open(path, "rb") as f:
            file_bytes = f.read()
        text, method = lire_pdf_robuste(file_bytes)
        return text, None, {"extraction_method": method}
    except Exception as e:
        return None, None, {"error": str(e)}

def parse_image_file(path):
    try:
        with open(path, "rb") as f:
            file_bytes = f.read()
        text, method = lire_image_directe(file_bytes)
        return text, {"extraction_method": method}
    except Exception:
        return None, None

def parse_msg_file(path):
    """ Lit les fichiers Outlook .msg """
    try:
        msg = extract_msg.Message(path)
        
        header = f"--- EMAIL OUTLOOK ---\nSujet: {msg.subject}\nDe: {msg.sender}\nPour: {msg.to}\nDate: {msg.date}\n"
        body = f"\n--- CORPS DU MESSAGE ---\n{msg.body}\n"
        
        return header + body, {"sender": msg.sender, "subject": msg.subject}
    except Exception as e:
        print(f"Erreur MSG: {e}")
        return None, None

def parse_pptx_file(path):
    try:
        prs = Presentation(path)
        text_content = []
        slide_count = 0
        for slide in prs.slides:
            slide_count += 1
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            text_content.append(f"--- SLIDE {slide_count} ---\n" + "\n".join(slide_text))
        return "\n\n".join(text_content), slide_count
    except Exception as e:
        return None, None

def parse_xlsx_file(path):
    try:
        all_sheets = pd.read_excel(path, sheet_name=None, dtype=str)
        
        full_text = ""
        sheets_metadata = []
        
        for sheet_name, df in all_sheets.items():
            df = df.dropna(how='all', axis=0).dropna(how='all', axis=1)
            df = df.fillna("")
            
            table_markdown = df.to_markdown(index=False)
            
            full_text += f"\n--- FEUILLE : {sheet_name} ---\n"
            full_text += table_markdown
            full_text += "\n\n"
            
            sheets_metadata.append(sheet_name)
            
        return full_text, {"sheets": sheets_metadata}
        
    except Exception as e:
        print(f"Erreur Excel Pandas: {e}")
        return None, None
    

def parse_pages(path):
    try:
        # On ouvre le .pages comme un fichier ZIP
        with zipfile.ZipFile(path, 'r') as z:
            file_list = z.namelist()
            
            # Apple range le PDF de prévisualisation à des endroits différents selon les versions
            preview_files = ['Preview.pdf', 'preview.pdf', 'QuickLook/Preview.pdf']
            
            for preview_file in preview_files:
                if preview_file in file_list:
                    # On a trouvé le PDF caché ! On l'extrait temporairement
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
                        with z.open(preview_file) as pdf_file:
                            temp_pdf.write(pdf_file.read())
                        temp_path = temp_pdf.name
                    
                    # On utilise notre super fonction PDF pour le lire
                    text, nb_pages, metadata = parse_pdf_robuste(temp_path)
                    
                    # On nettoie le fichier temporaire
                    os.unlink(temp_path)
                    
                    if metadata is None: metadata = {}
                    metadata["source"] = "pages_preview"
                    return text, nb_pages, metadata
            
            return None, None, {"error": "Aucun aperçu PDF trouvé (Sauvegarder avec aperçu dans Pages)"}
    except Exception as e:
        return None, None, {"error": f"Erreur .pages: {str(e)}"}    

def parse_docx_file(path):
    try:
        doc = docx.Document(path)
        full_text = []
        # Paragraphes
        for p in doc.paragraphs:
            full_text.append(p.text)

        for table in doc.tables:
            full_text.append("\n[TABLEAU WORD]")
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                full_text.append(" | ".join(row_text))
            full_text.append("[FIN TABLEAU]\n")
            
        return "\n".join(full_text), None
    except:
        return None, None

def parse_txt(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            lines = f.readlines()
        return "".join(lines), len(lines)
    except:
        return None, None

def parse_csv_file(path):
    try:
        rows = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        full_text = "\n".join([";".join(r) for r in rows])
        metadata = {"rows": len(rows)}
        return full_text, len(rows), metadata
    except:
        return None, None, None

def parse_json_file(path):
    try:
        data = json.load(open(path, "r"))
        pretty = json.dumps(data, indent=2, ensure_ascii=False)
        return pretty, {"type": str(type(data))}
    except:
        return None, None
    
def parse_pages(path):
    try:
        with zipfile.ZipFile(path, 'r') as z:
            file_list = z.namelist()
            preview_files = ['Preview.pdf', 'preview.pdf', 'QuickLook/Preview.pdf']
            
            for preview_file in preview_files:
                if preview_file in file_list:
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
                        with z.open(preview_file) as pdf_file:
                            temp_pdf.write(pdf_file.read())
                        temp_path = temp_pdf.name
                    
                    text, nb_pages, metadata = parse_pdf_robuste_wrapper(temp_path)
                    os.unlink(temp_path)
                    
                    if metadata is None: metadata = {}
                    metadata["source"] = "pages_preview"
                    return text, nb_pages, metadata
            return None, None, {"error": "Aucun aperçu PDF trouvé"}
    except Exception as e:
        return None, None, {"error": f"Erreur .pages: {str(e)}"}

def analyze_file(path, filename):
    size = os.path.getsize(path)
    if size > MAX_SIZE_BYTES:
        return None, "too_large"

    ext = filename.split(".")[-1].lower()
    mime = detect_mime(path)

    text, nb_pages, nb_lines, metadata = None, None, None, None

    if ext == "pdf":
        text, nb_pages, metadata = parse_pdf_robuste_wrapper(path)
    elif ext in ("png", "jpg", "jpeg", "tiff", "bmp"):
        text, metadata = parse_image_file(path)
        nb_pages = 1
    elif ext == "docx":
        text, nb_pages = parse_docx_file(path)
    elif ext == "pptx":
        text, nb_pages = parse_pptx_file(path)
    elif ext == "xlsx":
        text, metadata = parse_xlsx_file(path)
        nb_pages = 1
    elif ext == "msg":
        text, metadata = parse_msg_file(path)
        nb_pages = 1
    elif ext == "pages":
        text, nb_pages, metadata = parse_pages(path)  
    elif ext in ("txt", "md", "log"):
        text, nb_lines = parse_txt(path)
    elif ext == "csv":
        text, nb_lines, metadata = parse_csv_file(path)
    elif ext == "json":
        text, metadata = parse_json_file(path)
    else:
        text, nb_lines = parse_txt(path)
        if text: metadata = {"note": "format_inconnu_lu_comme_texte"}

    status = "parsed_ok" if text is not None else "parse_error"

    return {
        "name": filename,
        "extension": ext,
        "mime_type": mime,
        "size_bytes": size,
        "text_content": text,
        "nb_pages": nb_pages,
        "nb_lines": nb_lines,
        "metadata": json.dumps(metadata) if metadata else None,
        "status": status
    }, None